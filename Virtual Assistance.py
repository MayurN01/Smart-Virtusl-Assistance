from MySQLdb import Time
import pyttsx3
import speech_recognition as sr
import webbrowser
from pywikihow import search_wikihow
from bs4 import BeautifulSoup
import pywhatkit
import wikipedia
from googletrans import Translator
import os
import pyautogui
import psutil
from tkinter import Label
from tkinter import Entry
from tkinter import Button
import requests
from tkinter import Tk
from gtts import gTTS
from tkinter import StringVar
import PyPDF2
from pytube import YouTube
import datetime
from playsound import playsound
import keyboard
import pyjokes
from pptx.util import Inches, Pt
from pptx import Presentation

engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
engine.setProperty('voice', voices[0].id)
engine.setProperty('rate',170)

def Speak(Audio):
    print("   ")
    print(f": {Audio}")
    engine.say(Audio)
    print("    ")
    engine.runAndWait()
    

def takecommand(): 
    command = sr.Recognizer()
    with sr.Microphone(device_index=0) as source:
        print("Listening.......")
        command.pause_threshold = 1
        audio = command.listen(source, None, 5)

    try:
        print("Recognizing...")    
        query = command.recognize_google(audio, language='en-in')
        print(f"Your Command :  {query}\n")

    except:   
        return "None"
        
    return query.lower()


def TaskExe():
    Speak("Hello, I an Robert")
    Speak("How can I help you")

    def Music():
        Speak("Tell Me The NamE oF The Song!")
        musicName = takecommand()

        if 'akeli' in musicName:
            os.startfile('E:\\Songs\\akeli.mp3')

        elif 'blanko' in musicName:
            os.startfile('E:\\Songs\\blanko.mp3')

        else:
            pywhatkit.playonyt(musicName)

        Speak("Your Song Has Been Started! , Enjoy Sir!")

    def OpenApps():
        Speak("Ok Sir , Wait A Second!")
        
        if 'open vs code' in query:
            Speak("opening VS Code....")
            os.startfile("C:\\Users\\mayur\\AppData\\Local\\Programs\\Microsoft VS Code\\Code.exe")

        elif 'telegram' in query:
            Speak("opening Teleggram....")
            os.startfile("C:\\Users\\mayur\\OneDrive\\Desktop\\Telegram Desktop\\Telegram.exe")

        elif 'chrome' in query:
            Speak("opening Chrome....")
            os.startfile("C:\\Program Files\\Google\\Chrome\\Application\\chrome.exe")
         
        elif 'facebook' in query:
            Speak("opening Facebook....")
            webbrowser.open('https://www.facebook.com/')

        elif 'instagram' in query:
            Speak("opening Instagram....")
            webbrowser.open('https://www.instagram.com/')

        elif 'maps' in query:
            Speak("opening Google Maps....")
            webbrowser.open('https://www.google.com/maps/@28.7091225,77.2749958,15z')

        elif 'youtube' in query:
            Speak("opening Youtube....")
            webbrowser.open('https://www.youtube.com')
            
        elif 'open powerpoint' in query:
            Speak("opening powerpoint....")
            codePath = "C:\\Program Files\\Microsoft Office\\root\\Office16\\POWERPNT.EXE"
            os.startfile(codePath)
            
        elif 'open excel' in query:
            Speak("opening excel....")
            codePath = "C:\\Program Files\\Microsoft Office\\root\\Office16\\EXCEL.EXE"
            os.startfile(codePath)
            
        elif 'open wordpad' in query:
            Speak("opening wordpad....")
            codePath =  "C:\\Program Files\\Microsoft Office\\root\\Office16\\WINWORD.EXE"
            os.startfile(codePath)
            
        elif 'open notepad' in query:
            Speak("opening notepad....")
            codePath = "C:\\Windows\\notepad"
            os.startfile(codePath)

        elif 'open microsoft edge' in query:
            Speak("opening Microsoft Edge....")
            codePath = "C:\Program Files (x86)\Microsoft\Edge\Application\\msedge.exe"
            os.startfile(codePath)
            
        elif "write a note" in query:
            Speak("What should i write, sir")
            note = takecommand()
            file = open('jarvis.txt', 'w')
            Speak("Sir, Should i include date and time")
            snfm = takecommand()
            if 'yes' in snfm or 'sure' in snfm:
                strTime = datetime.datetime.now().strftime("%I:%M %p")
                file.write(strTime)
                file.write(" :- ")
                file.write(note)
            else:
                file.write(note)
            
            Speak("Note has been Saved")
                
        elif "show note" in query:
            Speak("Showing Notes")
            file = open("jarvis.txt", "r")
            print(file.read())
            Speak(file.read(6))
            
        
        elif "create powerpoint presentation" in query:
            root = Presentation () 
            first_layer = root.slide_layouts[0]  
            slide = root.slides.add_slide(first_layer) 
            Speak("Sir, what title should I add?")
            title = takecommand().upper()
            slide.shapes.title.text = title
          
            Speak("Sir, what subtitle should I add?")
            subtitle = takecommand().capitalize()
            slide.placeholders[1].text = subtitle
          
            Speak("sir, you want to add more slides?")
            command = takecommand()
            while "no" not in command or "yes" in command:
                second_slide_layout = root.slide_layouts[6]  #6 for blank slide
                slide = root.slides.add_slide(second_slide_layout)
                Speak("New slide added")
                Speak("What would you like to add in it?")
                print("Text\nImage\nTable\n")
                command = takecommand()
                if "text" in command:
                    left = top = width = height = Inches(1) 
                    txBox = slide.shapes.add_textbox(left, top,width, height) #creating the textbox
                    tf = txBox.text_frame   #creating the text frame
                    command = "yes"
                    while "no" not in command:
                        Speak("What Text you want to add")
                        p = tf.add_paragraph()
                        p.text = takecommand().capitalize()
                        Speak("Would you like to change font size")
                        command = takecommand().lower()
                        if "yes" in command:
                            Speak("Tell me the new font size")
                            command = int(takecommand())
                            p.font.size = Pt(command)
                        else:
                            p.font.size = Pt(20)

                        Speak("Would you like to bold your font or would you like to change your font style?")
                        command = takecommand()
                        if "yes" in command:
                            Speak("font or style??")
                            command = takecommand()
                            p.font.Calibri = True
                            if "style" in command:
                                print("italic\nVerdana\nCalibri\nPalatino\nTahoma\nGeorgia\nGill Sans\nCorbel\nSegoe\n")
                                Speak("Choose the style...")
                                command = takecommand()
                                p.font.command = True
                        
                            elif "font" in command:
                               p.font.bold = True
                            
                        Speak("Want to add more text in it?")
                        command = takecommand().lower()
                
                elif "image" in command:
                   # slide = root.slides.add_slide(second_slide_layout)
                    left = top = Inches(1)
                    height = Inches(5)
                    Speak("Name of the image")
                    img_path = takecommand()
                    pic = slide.shapes.add_picture(img_path+".jpg",left, top,height=height)
                    Speak("Image added")
                
                elif "table" in command:
                    x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
                    shape = slide.shapes.add_table(3, 4, x, y, cx, cy)
                    Speak("Table added")
                
                Speak("Would you want to add more slides?")
                command = takecommand()
            
            Speak("Tell me the name for this file?")
            file_name = takecommand()
            root.save(file_name+".pptx")
            Speak("File is saved")
            
            
        elif "open my presentation file" in query:
            Speak("What is the name of the file")
            file_name = takecommand()
            Speak("opening"+file_name+"....")
            codePath = file_name+".pptx"
            os.startfile(codePath)

        Speak("Your Command Has Been Completed Sir!")
        

    def Temp():
        
        search = "temperature in Nashik"
        url = f"https://www.google.com/search?q={search}"
        r = requests.get(url)
        data = BeautifulSoup(r.text,"html.parser")
        temperature = data.find("div",class_ = "BNeawe").text
        Speak(f"The Temperature Outside Is {temperature} ")

        Speak("Do I Have To Tell You Another Place Temperature ?")
        next = takecommand()

        if 'yes' in next:
            Speak("Tell Me The Name Of tHE Place ")
            name = takecommand()
            search = f"temperature in {name}"
            url = f"https://www.google.com/search?q={search}"
            r = requests.get(url)
            data = BeautifulSoup(r.text,"html.parser")
            temperature = data.find("div",class_ = "BNeawe").text
            Speak(f"The Temperature in {name} is {temperature} ")

        else:
            Speak("no problem sir")

    def Reader():
        Speak("Tell Me The Name Of The Book!")

        name = takecommand()

        if 'india' in name:

            os.startfile('E:\\Kaushik Shresth\\Books\\Social Science\\History\\ch 1.pdf')
            book = open('E:\\Kaushik Shresth\\Books\\Social Science\\History\\ch 1.pdf','rb')
            pdfreader = PyPDF2.PdfFileReader(book)
            pages = pdfreader.getNumPages()
            Speak(f"Number Of Pages In This Books Are {pages}")
            Speak("From Which Page I Have To Start Reading ?")
            numPage = int(input("Enter The Page Number :"))
            page = pdfreader.getPage(numPage)
            text = page.extractText()
            Speak("In Which Language , I Have To Read ?")
            lang = takecommand()

            if 'hindi' in lang:
                transl = Translator()
                textHin = transl.translate(text,'hi')
                textm = textHin.text
                speech = gTTS(text = textm )
                try:
                    speech.save('book.mp3')
                    playsound('book.mp3')

                except:
                    playsound('book.mp3')

            else:
                Speak(text)

        elif 'europe' in name:
            os.startfile('E:\\Kaushik Shresth\\Books\\Social Science\\History\\ch 3.pdf')
            book = open('E:\\Kaushik Shresth\\Books\\Social Science\\History\\ch 3.pdf','rb')
            pdfreader = PyPDF2.PdfFileReader(book)
            pages = pdfreader.getNumPages()
            Speak(f"Number Of Pages In This Books Are {pages}")
            Speak("From Which Page I Have To Start Reading ?")
            numPage = int(input())
            page = pdfreader.getPage(numPage)
            text = page.extractText()
            Speak("In Which Language , I Have To Read ?")
            lang = takecommand()

            if 'hindi' in lang:
                transl = Translator()
                textHin = transl.translate(text,'hi')
                textm = textHin.text
                speech = gTTS(text = textm )
                try:

                    speech.save('book.mp3')
                    playsound('book.mp3')

                except:
                    playsound('book.mp3')

            else:
                Speak(text)

    def CloseAPPS():
        Speak("Ok Sir , Wait A second!")

        if 'youtube' in query:
            os.system("TASKKILL /F /im Chrome.exe")

        elif 'chrome' in query:
            os.system("TASKKILL /F /IM Chrome.exe")

        elif 'telegram' in query:
            os.system("TASKKILL /F /im Telegram.exe")

        elif 'code' in query:
            os.system("TASKKILL /F /im code.exe")

        elif 'instagram' in query:
            os.system("TASKKILL /F /im chrome.exe")
            
            
        Speak("Your Command Has Been Succesfully Completed!")

    def YoutubeAuto():
        Speak("Whats Your Command ?")
        comm = takecommand()

        if 'pause' in comm:
            keyboard.press('space bar')

        elif 'restart' in comm:
            keyboard.press('0')

        elif 'mute' in comm:
            keyboard.press('m')

        elif 'skip' in comm:
            keyboard.press('l')

        elif 'back' in comm:
            keyboard.press('j')

        elif 'full screen' in comm:
            keyboard.press('f')

        elif 'film mode' in comm:
            keyboard.press('t')

        Speak("Done Sir")

    def TakeHindi():
        command = sr.Recognizer()
        with sr.Microphone() as source:
            print("Listening......")
            command.pause_threshold = 1
            audio = command.listen(source)

            try:
                print("Recognizing.....")
                query = command.recognize_google(audio,language='hi')
                print(f"You Said : {query}")

            except:
                return "none"

            return query.lower()

    def Tran():
        Speak("Tell Me The Line!")
        line = TakeHindi()
        traslate = Translator()
        result = traslate.translate(line)
        Text = result.text
        Speak(Text)
        
    def ChromeAuto():
        Speak("Chrome Automation started!")

        command = takecommand()

        if 'close this tab' in command:
            keyboard.press_and_release('ctrl + w')

        elif 'open new tab' in command:
            keyboard.press_and_release('ctrl + t')

        elif 'open new window' in command:
            keyboard.press_and_release('ctrl + n')

        elif 'history' in command:
            keyboard.press_and_release('ctrl +h')

    def screenshot():
        Speak("Ok Boss , What Should I Name That File ?")
        path = takecommand()
        path1name = path + ".png"
        path1 = "D:\\VS Code Files\\"+ path1name
        kk = pyautogui.screenshot()
        kk.save(path1)
        os.startfile("D:\\VS Code Files")
        Speak("Here Is Your ScreenShot") 

    while True:
    
        query = takecommand()

        if 'hello' in query:
            Speak("Hello Sir , I Am Robert .")
            Speak("Your Personal AI Assistant!")
            Speak("How May I Help You?")

        elif 'how are you' in query:
            Speak("I Am Fine Sir!")
            Speak("Whats About YOU?")

        elif 'take a break' in query:
            Speak("Ok Sir , You Can Call Me Anytime !")
            Speak("Just Say Wake Up Robert!")
            break

        elif 'youtube search' in query:
            Speak("OK sIR , This Is What I found For Your Search!")
            query = query.replace("Robert","")
            query = query.replace("youtube search","")
            web = 'https://www.youtube.com/results?search_query=' + query
            webbrowser.open(web)
            Speak("Done Sir!")

        elif 'website' in query:
            Speak("Ok Sir , Launching.....")
            query = query.replace("Robert","")
            query = query.replace("website","")
            query = query.replace(" ","")
            web1 = query.replace("open","")
            web2 = 'https://www.' + web1 + '.com'
            webbrowser.open(web2)
            Speak("Launched!")

        elif 'launch' in query:
            Speak("Tell Me The Name Of The Website!")
            name = takecommand()
            web = 'https://www.' + name + '.com'
            webbrowser.open(web)
            Speak("Done Sir!")

        elif 'wikipedia' in query:
            Speak("Searching Wikipedia.....")
            query = query.replace("Robert","")
            query = query.replace("wikipedia","")
            wiki = wikipedia.summary(query,2)
            Speak(f"According To Wikipedia : {wiki}")
            
        elif 'the time' in query:
            strTime = datetime.datetime.now().strftime("%I:%M %p")
            Speak(f"Sir, the time is {strTime}")
        
        elif 'screenshot' in query:
            screenshot()

        elif 'open facebook' in query:
            OpenApps()

        elif 'open instagram' in query:
            OpenApps()

        elif 'open maps' in query:
            OpenApps()

        elif 'open vs code' in query:
            OpenApps()

        elif 'open youtube' in query:
            OpenApps()
            
        elif 'open telegram' in query:
            OpenApps()

        elif 'open chrome' in query:
            OpenApps()
            
        elif 'open excel' in query:
            OpenApps()
            
        elif 'open notepad' in query:
            OpenApps()
            
        elif 'open powerpoint' in query:
            OpenApps()
            
        elif 'open wordpad' in query:
            OpenApps()
            
        elif 'open microsoft edge' in query:
            OpenApps()
            
        elif 'write a note' in query:
            OpenApps()
            
        elif 'show note' in query:
            OpenApps()
            
        elif 'create powerpoint presentation' in query:
            OpenApps()
            
        elif 'open my presentation file' in query:
            OpenApps()

        elif 'close chrome' in query:
            CloseAPPS()

        elif 'music' in query:
            Music()

        elif 'close telegram' in query:
            CloseAPPS()

        elif 'close instagram' in query:
            CloseAPPS()

        elif 'close facebook' in query:
            CloseAPPS()

        elif 'pause' in query:
            keyboard.press('space bar')

        elif 'restart' in query:
            keyboard.press('0')

        elif 'mute' in query:
            keyboard.press('m')

        elif 'skip' in query:
            keyboard.press('l')

        elif 'back' in query:
            keyboard.press('j')

        elif 'full screen' in query:
            keyboard.press('f')

        elif 'film mode' in query:
            keyboard.press('t')

        elif 'start youtube ' in query:
            YoutubeAuto()

        elif 'close the tab' in query:
            keyboard.press_and_release('ctrl + w')

        elif 'open new tab' in query:
            keyboard.press_and_release('ctrl + t')

        elif 'open new window' in query:
            keyboard.press_and_release('ctrl + n')

        elif 'history' in query:
            keyboard.press_and_release('ctrl +h')

        elif 'start chrome automation' in query:
            ChromeAuto()

        elif 'tell me a joke' in query:
            get = pyjokes.get_joke()
            Speak(get)

        elif 'repeat my word' in query:
            Speak("Speak Sir!")
            jj = takecommand()
            Speak(f"You Said : {jj}")

        elif 'my location' in query:
            Speak("Ok Sir , Wait A Second!")
            webbrowser.open('https://www.google.com/maps/@28.7091225,77.2749958,15z')

        elif 'set a alarm' in query:
            Speak("Enter The Time !")
            time = input(": Enter The Time :")

            while True:
                Time_Ac = datetime.datetime.now()
                now = Time_Ac.strftime("%H:%M:%S")

                if now == time:
                    Speak("Time To Wake Up Sir!")
                    playsound('iron.mp3')
                    Speak("Alarm Closed!")

                elif now>time:
                    break

        elif 'video downloader' in query:
            root = Tk()
            root.geometry('500x300')
            root.resizable(0,0)
            root.title("Youtube Video Downloader")
            Speak("Enter Video Url Here !")
            Label(root,text = "Youtube Video Downloader",font = 'arial 15 bold').pack()
            link = StringVar()
            Label(root,text = "Paste Yt Video URL Here",font = 'arial 15 bold').place(x=160,y=60)
            Entry(root,width = 70,textvariable = link).place(x=32,y=90)

            def VideoDownloader():
                url = YouTube(str(link.get()))
                video = url.streams.first()
                video.download()
                Label(root,text = "Downloaded",font = 'arial 15').place(x= 180,y=210)

            Button(root,text = "Download",font = 'arial 15 bold',bg = 'pale violet red',padx = 2 , command = VideoDownloader).place(x=180,y=150)

            root.mainloop()
            Speak("Video Downloaded")
            
        elif 'translator' in query:
            Tran()
    
        elif 'Talk in hindi' in query:
            TakeHindi()
        
        elif 'remember that' in query:
            remeberMsg = query.replace("remember that","")
            remeberMsg = remeberMsg.replace("Robert","")
            Speak("You Tell Me To Remind You That :"+remeberMsg)
            remeber = open('data.txt','w')
            remeber.write(remeberMsg)
            remeber.close()

        elif 'what do you remember' in query:
            remeber = open('data.txt','r')
            Speak("You Tell Me That" + remeber.read())

        elif 'google search' in query:
            import wikipedia as googleScrap
            query = query.replace("Robert","")
            query = query.replace("google search","")
            query = query.replace("google","")
            Speak("This Is What I Found On The Web!")
            pywhatkit.search(query)

            try:
                result = googleScrap.summary(query,2)
                Speak(result)

            except:
                Speak("No Speakable Data Available!")

        elif 'how to' in query:
            Speak("Getting Data From The Internet !")
            op = query.replace("Robert","")
            max_result = 1
            how_to_func = search_wikihow(op,max_result)
            assert len(how_to_func) == 1
            how_to_func[0].print()
            Speak(how_to_func[0].summary)
            
        elif 'the temperature' in query:
            Temp()

        elif 'read book' in query:
            Reader()
            
        else:
            Speak("I cant get that")          
    

TaskExe()